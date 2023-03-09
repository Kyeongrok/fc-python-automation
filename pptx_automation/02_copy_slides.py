import copy

from pptx import Presentation


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide0 = self.ppt_file.slides[slide_no]

        for shape in slide0.shapes:
            print(f"slide_no:{slide_no} {shape.text}")

    def copy_slide(self, from_slide_no=0, slide_layout_no=6):
        from_slide = self.ppt_file.slides[from_slide_no]

        to_slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[slide_layout_no])

        for shape in from_slide.shapes:
            el = shape.element
            new_element = copy.deepcopy(el)
            to_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    def duplicate_n_slides(self, slide_cnt, from_slide_no=0):
        # 1000개를 원한다면 999개를 넘겨주세요. 기존에 1개 있기 때문입니다.
        for _ in range(slide_cnt):
            self.copy_slide(from_slide_no=from_slide_no)

    def save(self, filename):
        self.ppt_file.save(filename)

    def get_shape_map(self, slide_no):
        slide = self.ppt_file.slides[slide_no]
        shape_map = {}
        for i, shape in enumerate(slide.shapes):
            shape_map[shape.name] = i
        return shape_map

    def change_text(self, slide_no, label_map):
        slide = self.ppt_file.slides[slide_no]
        shape_map = self.get_shape_map(slide_no)

        for shape_name, text in label_map.items():
            shape_no = shape_map[shape_name]
            slide.shapes[shape_no].text = text


if __name__ == '__main__':
    ppt_al = PowerPointAutoLabel("재물조사표.pptx")
    # ppt_al.print_slide_shapes(0)
    ppt_al.duplicate_n_slides(1)
    ppt_al.change_text(0, {"product_name":"32인치 게이밍 모니터", "model_no":"MO11223344"})
    ppt_al.change_text(1, {"product_name":"40인치 거실용 TV", "model_no":"MO7890"})
    ppt_al.print_slide_shapes(0)
    ppt_al.save("[Auto]재물조사표_change_text.pptx")
